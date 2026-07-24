#!/usr/bin/env python3
"""Phase-10: build the first six slides of the performance-gap deck as a NEW draft.

The user's master deck (`performance-gap-experiments-presentation.pptx`) is not in
this repository, so per the plan we create a *draft* rather than overwriting a
source file. Style follows the existing light deck: white background, dark-navy
headings, muted blue/green accents, few words, large charts, speaker notes.
"""
from __future__ import annotations
import argparse
import json
from pathlib import Path

import pandas as pd
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.util import Emu, Inches, Pt

NAVY = RGBColor(0x1B, 0x2A, 0x49)
BLUE = RGBColor(0x3D, 0x6F, 0xB4)
GREEN = RGBColor(0x3F, 0x8F, 0x6B)
ORANGE = RGBColor(0xD0, 0x84, 0x28)
GREY = RGBColor(0x5A, 0x63, 0x74)
W, H = Inches(13.333), Inches(7.5)


def blank(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    return s


def title(slide, text, sub=None):
    tb = slide.shapes.add_textbox(Inches(0.55), Inches(0.32), W - Inches(1.1), Inches(0.85))
    tf = tb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.text = text
    p.runs[0].font.size = Pt(30); p.runs[0].font.bold = True
    p.runs[0].font.color.rgb = NAVY
    if sub:
        p2 = tf.add_paragraph(); p2.text = sub
        p2.runs[0].font.size = Pt(13.5); p2.runs[0].font.color.rgb = GREY
    return tb


def bullets(slide, items, left, top, width, height, size=15, color=NAVY,
            head_color=BLUE):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame; tf.word_wrap = True
    first = True
    for it in items:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        if isinstance(it, tuple):          # (heading, None) style
            p.text = it[0]
            p.runs[0].font.size = Pt(size + 1.5); p.runs[0].font.bold = True
            p.runs[0].font.color.rgb = head_color
            p.space_before = Pt(9)
        else:
            p.text = "•  " + it
            p.runs[0].font.size = Pt(size); p.runs[0].font.color.rgb = color
            p.space_before = Pt(3)
    return tb


def pic(slide, path, left, top, width=None, height=None):
    if not Path(path).exists():
        return None
    return slide.shapes.add_picture(str(path), left, top, width=width, height=height)


def footer(slide, text):
    tb = slide.shapes.add_textbox(Inches(0.55), H - Inches(0.42), W - Inches(1.1),
                                  Inches(0.32))
    p = tb.text_frame.paragraphs[0]; p.text = text
    p.runs[0].font.size = Pt(9); p.runs[0].font.color.rgb = GREY


def notes(slide, text):
    slide.notes_slide.notes_text_frame.text = text


def banner(slide, text, top, color=GREEN):
    from pptx.enum.shapes import MSO_SHAPE
    sh = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.55), top,
                                W - Inches(1.1), Inches(0.62))
    sh.fill.solid(); sh.fill.fore_color.rgb = RGBColor(0xF2, 0xF6, 0xF3)
    sh.line.color.rgb = color; sh.line.width = Pt(1.2)
    tf = sh.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.text = text
    p.runs[0].font.size = Pt(14); p.runs[0].font.bold = True
    p.runs[0].font.color.rgb = NAVY
    return sh


def build(root: Path, out: Path):
    sm = pd.read_csv(root / "summary_matrix.csv")
    env = json.loads((root / "environment.json").read_text())
    plots = root / "plots"
    models = sorted(sm.model.unique())
    n_cfg = int(sm.n_configs_evaluated.max())
    src = "results/2026-07-24_serving_ceiling"

    prs = Presentation(); prs.slide_width, prs.slide_height = W, H

    # ---------------- Slide 1 ----------------
    s = blank(prs)
    title(s, "Finding the Performance Gap in LLM Inference",
          "Qwen3-30B-A3B and LFM2.5-8B-A1B · single H200 · TP1 · BF16 · SGLang")
    bullets(s, [
        ("The question",),
        "A model that runs is not a model that runs well — how much of the gap is "
        "just serving configuration?",
        ("What we do first",),
        "Freeze a fair cookbook baseline, then enumerate the serving-configuration "
        f"space ({n_cfg} configurations) across six workload regimes.",
        "Measure throughput and TTFT / TPOT / E2E percentiles from raw per-request "
        "records — no metric is estimated from another.",
        ("Why it matters",),
        "Whatever serving tuning cannot close is the gap that profiling, kernels, "
        "backends and communication must explain.",
    ], Inches(0.6), Inches(1.75), W - Inches(1.3), Inches(4.6), size=16)
    banner(s, "Goal: quantify the serving-level ceiling before blaming the kernels.",
           Inches(6.15))
    footer(s, f"{src} · sglang {env['sglang_version']} @ {env['sglang_commit'][:9]} · "
              f"torch {env['torch']} · triton {env['triton']} · driver {env['driver']}")
    notes(s, "Framing slide. The deliverable of this section is an upper bound on "
             "what serving-level configuration alone can buy, measured under one "
             "frozen software stack. Everything after this slide is measured, not "
             "estimated.")

    # ---------------- Slide 2 ----------------
    s = blank(prs)
    title(s, "Runnable does not mean workload-optimal")
    bullets(s, [
        ("Runnable deployment",),
        "The official cookbook command starts the model and serves traffic correctly.",
        ("Hidden performance gap",),
        "The same command is a fixed point in a large configuration space; nothing "
        "in it adapts to the traffic you actually receive.",
        ("The gap moves with the workload",),
        "Admission capacity, prefill chunking and scheduling interact differently "
        "for short decode, long prefill, high concurrency and agentic traffic.",
    ], Inches(0.6), Inches(1.8), Inches(6.1), Inches(4.4), size=15.5)
    p = pic(s, plots / "gain_distribution.png", Inches(6.95), Inches(1.85),
            width=Inches(5.9))
    banner(s, "Same deployment, different traffic — the reachable gain is not a "
              "single number.", Inches(6.2))
    footer(s, f"{src}/plots/gain_distribution.png")
    notes(s, "Problem definition. The right-hand chart already answers 'is tuning "
             "worth it?' with a distribution rather than a headline: for each regime "
             "it shows what fraction of the evaluated configurations are wins, "
             "flats, regressions or trade-offs relative to that model's own cookbook.")

    # ---------------- Slide 3 ----------------
    s = blank(prs)
    title(s, "Freeze a fair baseline before measuring the gap")
    bullets(s, [
        ("What the cookbook gives",),
        "Official deployment command for this model / GPU / dtype cell.",
        "A supported, reproducible starting configuration.",
        ("What we additionally freeze",),
        f"Software: sglang {env['sglang_version']} @ {env['sglang_commit'][:9]}, "
        f"torch {env['torch']}, triton {env['triton']}.",
        "Resolved execution path: attention backend fa3, MoE runner auto, "
        "CUDA Graph on — parsed from every server log, not assumed.",
        "Identical workloads, seeds and request payloads for every configuration.",
        ("Why it matters",),
        "No stale or handicapped baseline: the 2026-06-11 '5–9×' result came from a "
        "CUDA-Graph-disabled reference and collapses to 1.00–1.05× against the true "
        "default.",
        "Every claimed gap has a reproducible reference point.",
    ], Inches(0.6), Inches(1.75), W - Inches(1.3), Inches(4.7), size=14.5)
    banner(s, "The baseline is measured under the identical protocol — it is "
              "configuration #74 inside the same grid.", Inches(6.35), ORANGE)
    footer(s, f"{src}/baseline_definition.json · "
              "docs/2026-06-25/autotuning_honest_results.md")
    notes(s, "Baseline-fairness slide. The key honesty point: our own earlier 5–9× "
             "number was an artifact of a CUDA-Graph-disabled baseline; measured "
             "against the true zero-flag default the same winner is 1.00–1.05×. We "
             "therefore run the cookbook config inside the grid under the identical "
             "harness instead of quoting a historical number.")

    # ---------------- Slide 4 ----------------
    s = blank(prs)
    title(s, "How we search the serving configuration space")
    pic(s, plots / "search_space_overview.png", Inches(0.6), Inches(1.5),
        width=W - Inches(1.2))
    footer(s, f"{src}/search_space.yaml · {src}/workloads.yaml")
    notes(s, "Method slide. Full grid enumeration (not TPE) because we need the "
             "whole Pareto frontier and the negative results, not one winner. No "
             "warm start, no seeded cookbook, no reused study. Backends and CUDA "
             "Graph are frozen so that anything we measure is attributable to the "
             "four serving knobs. Note chunked_prefill_size=8192 is effectively "
             "equal to -1 here because context length is 8192 — that pair doubles "
             "as a built-in noise estimate.")

    # ---------------- Slide 5 ----------------
    s = blank(prs)
    title(s, "What serving tuning actually finds: wins, flats, regressions, trade-offs")
    pic(s, plots / "full_result_matrix_heatmap.png", Inches(0.5), Inches(1.35),
        width=Inches(8.5))
    lines = [("Read the matrix",)]
    for model in models:
        d = sm[sm.model == model]
        best = d.loc[d.d_request_throughput.idxmax()]
        worst = d.loc[d.d_request_throughput.idxmin()]
        lines.append(f"{model}: largest gain {best.d_request_throughput*100:+.0f}% "
                     f"({best.workload}); smallest {worst.d_request_throughput*100:+.0f}% "
                     f"({worst.workload}).")
    lines += [
        ("Also on the record",),
        "Historical Qwen autotuning vs a correct baseline: 1.00–1.05× (flat).",
        "High-concurrency stress study: throughput 1.40–2.44× and TTFT p50/p95 "
        "−85…−96 %, driven almost entirely by admission capacity.",
        "The same chunking candidate helps LFM2.5 shared-prefix (+28.6 % req/s) "
        "and is neutral on Qwen (−2.9 %).",
    ]
    bullets(s, lines, Inches(9.15), Inches(1.5), Inches(3.75), Inches(4.7), size=11.5)
    banner(s, "Serving tuning finds workload-specific cliffs, not a universal winner.",
           Inches(6.35))
    footer(s, f"{src}/summary_matrix.csv · "
              "results/2026-07-23_high_concurrency_ttft_rerun/comparison.md · "
              "results/consolidated_v7_config_sweep.csv")
    notes(s, "Result slide. Green cells are improvements, red are regressions; the "
             "label at the right of each row is the WIN/FLAT/REGRESSION/TRADE-OFF "
             "classification. Deliberately includes flat and negative regimes. The "
             "high-concurrency and cross-model numbers come from separate, clearly "
             "labelled studies with a different workload definition — they are "
             "supporting evidence, not cells of this matrix.")

    # ---------------- Slide 6 ----------------
    s = blank(prs)
    title(s, "Synthetic winners do not transfer reliably to agentic workloads")
    m0 = models[0]
    pic(s, plots / f"transfer_matrix_request_throughput_{m0}.png",
        Inches(0.5), Inches(1.45), width=Inches(6.5))
    pic(s, plots / f"per_regime_pareto_grid_{m0}.png",
        Inches(7.1), Inches(1.45), width=Inches(5.8))
    bullets(s, [
        "Rows are configurations that won some regime; columns are the regimes they "
        "are then applied to. Off-diagonal cells below 1.00× are failed transfers.",
        "Shared-prefix and tool-agent Pareto fronts show all evaluated points; "
        "cookbook is the star, the validated throughput winner is the diamond.",
    ], Inches(0.6), Inches(5.55), W - Inches(1.3), Inches(0.9), size=12.5)
    banner(s, "Serving search selects points on the frontier.  Profiling is needed "
              "to move the frontier.", Inches(6.5), BLUE)
    footer(s, f"{src}/analysis/{m0}/transfer_matrix_request_throughput.csv · "
              f"{src}/analysis/{m0}/pareto_points.csv")
    notes(s, "Transfer slide, and the hand-off into the profiling/kernel section. "
             "The diagonal is by construction >= 1.00x; what matters is how far the "
             "off-diagonal cells fall below it, especially synthetic -> agentic. "
             "Ratios are computed against each target regime's own cookbook, so the "
             "matrix is dimensionless and comparable across regimes.")

    out.parent.mkdir(parents=True, exist_ok=True)
    prs.save(out)
    print("wrote", out)


def main():
    ap = argparse.ArgumentParser()
    ap.add_argument("--outroot", required=True)
    ap.add_argument("--out", required=True)
    a = ap.parse_args()
    build(Path(a.outroot), Path(a.out))


if __name__ == "__main__":
    main()
